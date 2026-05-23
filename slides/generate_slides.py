from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ===== カラーパレット =====
C_DARK   = RGBColor(0x1A, 0x1A, 0x2E)   # 濃紺（背景）
C_MID    = RGBColor(0x16, 0x21, 0x3E)   # 中紺
C_ACCENT = RGBColor(0x0F, 0x3A, 0x5C)   # アクセント紺
C_BLUE   = RGBColor(0x00, 0x8B, 0xD4)   # 明るい青
C_CYAN   = RGBColor(0x00, 0xD4, 0xFF)   # シアン
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_LGRAY  = RGBColor(0xCC, 0xDD, 0xEE)
C_YELLOW = RGBColor(0xFF, 0xD7, 0x00)
C_GREEN  = RGBColor(0x00, 0xE6, 0x76)
C_RED    = RGBColor(0xFF, 0x4D, 0x4D)
C_ORANGE = RGBColor(0xFF, 0xA5, 0x00)

W = Inches(13.33)   # ワイドスクリーン幅
H = Inches(7.5)     # 高さ

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # 完全白紙レイアウト

# ===================================================
# ユーティリティ
# ===================================================
def add_rect(slide, l, t, w, h, fill=None, line=None, line_w=Pt(0)):
    from pptx.util import Pt
    shape = slide.shapes.add_shape(1, l, t, w, h)  # MSO_SHAPE_TYPE.RECTANGLE=1
    shape.line.width = line_w
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h,
             font_size=Pt(18), bold=False, color=C_WHITE,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(l, t, w, h)
    txb.word_wrap = wrap
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = font_size
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    run.font.name  = "Meiryo UI"
    return txb

def set_bg(slide, color=C_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_slide():
    sl = prs.slides.add_slide(BLANK)
    set_bg(sl)
    return sl

def add_header_bar(slide, title_text, sub_text=""):
    """上部タイトルバー"""
    add_rect(slide, Inches(0), Inches(0), W, Inches(1.15), fill=C_ACCENT)
    # アクセントライン
    add_rect(slide, Inches(0), Inches(1.10), W, Inches(0.07), fill=C_BLUE)
    add_text(slide, title_text,
             Inches(0.4), Inches(0.1), Inches(12), Inches(0.75),
             font_size=Pt(28), bold=True, color=C_WHITE)
    if sub_text:
        add_text(slide, sub_text,
                 Inches(0.4), Inches(0.82), Inches(12), Inches(0.35),
                 font_size=Pt(14), color=C_CYAN)

def add_card(slide, l, t, w, h, fill=C_MID, border=C_BLUE, bw=Pt(1.5)):
    add_rect(slide, l, t, w, h, fill=fill, line=border, line_w=bw)

def bullet_lines(slide, lines, l, t, w, h,
                 size=Pt(17), color=C_WHITE, spacing=Inches(0.38)):
    """箇条書き複数行"""
    for i, (bullet, text) in enumerate(lines):
        ty = t + i * spacing
        if bullet:
            add_text(slide, bullet, l, ty, Inches(0.35), Inches(0.38),
                     font_size=size, color=C_CYAN, bold=True)
            add_text(slide, text, l + Inches(0.38), ty, w - Inches(0.38), Inches(0.38),
                     font_size=size, color=color)
        else:
            add_text(slide, text, l, ty, w, Inches(0.38),
                     font_size=size, color=color)

def page_num(slide, n, total=14):
    add_text(slide, f"{n} / {total}",
             Inches(11.8), Inches(7.1), Inches(1.4), Inches(0.3),
             font_size=Pt(11), color=C_LGRAY, align=PP_ALIGN.RIGHT)

# ===================================================
# スライド 1: タイトル
# ===================================================
sl = add_slide()
# グラデーション風帯
add_rect(sl, Inches(0), Inches(2.3), W, Inches(3.1), fill=C_ACCENT)
add_rect(sl, Inches(0), Inches(2.3), W, Inches(0.06), fill=C_BLUE)
add_rect(sl, Inches(0), Inches(5.35), W, Inches(0.06), fill=C_BLUE)

add_text(sl, "RISC-Vプロセッサの高速化設計と評価",
         Inches(0.7), Inches(2.55), Inches(11.9), Inches(1.1),
         font_size=Pt(36), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
add_text(sl, "パイプライン処理による実行時間・消費電力の比較",
         Inches(0.7), Inches(3.65), Inches(11.9), Inches(0.7),
         font_size=Pt(22), color=C_CYAN, align=PP_ALIGN.CENTER)

add_text(sl, "コンピュータアーキテクチャ  グループ発表",
         Inches(0.7), Inches(5.6), Inches(11.9), Inches(0.45),
         font_size=Pt(16), color=C_LGRAY, align=PP_ALIGN.CENTER)
add_text(sl, "2025年",
         Inches(0.7), Inches(6.1), Inches(11.9), Inches(0.45),
         font_size=Pt(14), color=C_LGRAY, align=PP_ALIGN.CENTER)
page_num(sl, 1)

# ===================================================
# スライド 2: アジェンダ
# ===================================================
sl = add_slide()
add_header_bar(sl, "目次")
page_num(sl, 2)

items = [
    ("01", "目的"),
    ("02", "高速化設計の方針"),
    ("03", "ソートアルゴリズムの比較"),
    ("04", "高速化設計の詳細"),
    ("05", "評価結果"),
    ("06", "考察"),
    ("07", "感想・役割分担・エフォート"),
]
for i, (num, label) in enumerate(items):
    col = 0 if i < 4 else 1
    row = i if i < 4 else i - 4
    lx = Inches(1.0 + col * 6.3)
    ty = Inches(1.45 + row * 1.35) if i < 4 else Inches(1.45 + row * 1.35)
    add_card(sl, lx, ty, Inches(5.8), Inches(1.1))
    add_text(sl, num, lx + Inches(0.15), ty + Inches(0.12),
             Inches(0.7), Inches(0.55),
             font_size=Pt(22), bold=True, color=C_CYAN)
    add_text(sl, label, lx + Inches(0.8), ty + Inches(0.2),
             Inches(4.8), Inches(0.55),
             font_size=Pt(20), bold=True, color=C_WHITE)

# ===================================================
# スライド 3: 目的
# ===================================================
sl = add_slide()
add_header_bar(sl, "目的", "01 / Purpose")
page_num(sl, 3)

add_card(sl, Inches(0.5), Inches(1.4), Inches(12.3), Inches(2.1), fill=C_MID)
add_text(sl,
         "FPGA上にRISC-Vプロセッサを実装し，パイプライン処理を導入することで\n"
         "実行時間・消費電力・消費エネルギを改善することを目的とする．",
         Inches(0.9), Inches(1.55), Inches(11.5), Inches(1.8),
         font_size=Pt(20), color=C_WHITE)

lines_obj = [
    ("▶", "逐次実行CPU（選択ソート・クイックソート）を基準として比較"),
    ("▶", "2段パイプライン / 5段パイプラインの効果を定量的に評価"),
    ("▶", "実行時間・CPI・回路資源・消費電力・消費エネルギを指標とする"),
]
bullet_lines(sl, lines_obj,
             Inches(0.9), Inches(3.8), Inches(11.5), Inches(2.0),
             size=Pt(18), spacing=Inches(0.7))

# ===================================================
# スライド 4: 高速化設計の方針
# ===================================================
sl = add_slide()
add_header_bar(sl, "高速化設計の方針", "02 / Design Policy")
page_num(sl, 4)

cards = [
    ("アルゴリズム選択", "選択ソートとクイックソートを実装し，\n実行命令数の少ない方を採用"),
    ("動作周波数の向上", "クリティカルパスを短縮し\n高い周波数で動作させる"),
    ("パイプライン化", "命令実行を複数ステージに分割し\n複数命令を並列に処理"),
]
for i, (title, body) in enumerate(cards):
    lx = Inches(0.55 + i * 4.27)
    add_card(sl, lx, Inches(1.45), Inches(4.0), Inches(4.8))
    add_rect(sl, lx, Inches(1.45), Inches(4.0), Inches(0.55), fill=C_BLUE)
    add_text(sl, f"Step {i+1}  {title}",
             lx + Inches(0.15), Inches(1.47), Inches(3.7), Inches(0.5),
             font_size=Pt(15), bold=True, color=C_DARK)
    add_text(sl, body,
             lx + Inches(0.2), Inches(2.15), Inches(3.6), Inches(3.8),
             font_size=Pt(17), color=C_WHITE)

add_text(sl, "※ パイプライン段数を増やすとクロックあたりの処理が細分化され，高周波動作が可能になるが，\n　ハザードによるストールでCPIが増加するトレードオフが生じる",
         Inches(0.7), Inches(6.6), Inches(11.9), Inches(0.75),
         font_size=Pt(13), color=C_LGRAY, italic=True)

# ===================================================
# スライド 5: ソートアルゴリズムの比較
# ===================================================
sl = add_slide()
add_header_bar(sl, "ソートアルゴリズムの比較", "03 / Algorithm Comparison")
page_num(sl, 5)

# 表ヘッダ
cols = [Inches(1.5), Inches(4.3), Inches(7.6), Inches(10.5)]
col_w = [Inches(2.7), Inches(3.2), Inches(2.8), Inches(2.7)]
headers = ["指標", "選択ソート（逐次）", "クイックソート（逐次）", "優劣"]
add_rect(sl, Inches(0.5), Inches(1.45), Inches(12.3), Inches(0.5), fill=C_BLUE)
for j, h in enumerate(headers):
    add_text(sl, h, cols[j], Inches(1.47), col_w[j], Inches(0.45),
             font_size=Pt(15), bold=True, color=C_DARK, align=PP_ALIGN.CENTER)

rows = [
    ("実行命令数",    "3,429,518",  "118,777",   "クイック 約29倍少"),
    ("動作周波数",    "50 MHz",     "50 MHz",    "—"),
    ("実行時間",      "68.59 ms",   "2.376 ms",  "クイック 約29倍高速"),
    ("消費電力",      "0.394 W",    "0.384 W",   "ほぼ同等"),
    ("LUT使用率",     "29.90%",     "29.90%",    "同等"),
]
for i, (lbl, v1, v2, note) in enumerate(rows):
    bg = C_MID if i % 2 == 0 else C_ACCENT
    add_rect(sl, Inches(0.5), Inches(1.95 + i * 0.9), Inches(12.3), Inches(0.85), fill=bg)
    vals = [lbl, v1, v2, note]
    cs = [C_CYAN, C_WHITE, C_WHITE, C_YELLOW]
    for j in range(4):
        add_text(sl, vals[j], cols[j], Inches(1.97 + i * 0.9), col_w[j], Inches(0.8),
                 font_size=Pt(15), color=cs[j], align=PP_ALIGN.CENTER)

add_text(sl, "→ 実行命令数が約29倍少ないクイックソートを採用し，パイプライン化の対象とする",
         Inches(0.7), Inches(6.7), Inches(11.9), Inches(0.55),
         font_size=Pt(16), bold=True, color=C_GREEN)

# ===================================================
# スライド 6: 高速化設計の詳細（逐次CPU）
# ===================================================
sl = add_slide()
add_header_bar(sl, "高速化設計の詳細 ① 逐次実行CPU", "04 / Sequential CPU")
page_num(sl, 6)

add_card(sl, Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.8))
add_text(sl, "設計概要",
         Inches(0.8), Inches(1.5), Inches(11), Inches(0.5),
         font_size=Pt(20), bold=True, color=C_CYAN)

lines = [
    ("●", "RISC-V（RV32I）命令セットを実装したシングルサイクルCPU"),
    ("●", "1クロックで1命令を完了（CPI ≈ 1.0）"),
    ("●", "動作周波数：50 MHz（クリティカルパスが長いため制限あり）"),
    ("●", "命令メモリ・データメモリを内蔵（Distributed RAM使用）"),
    ("●", "クイックソートで 118,777 命令，実行時間 2.376 ms"),
]
bullet_lines(sl, lines, Inches(0.9), Inches(2.1), Inches(11.3), Inches(3.5),
             size=Pt(17), spacing=Inches(0.62))

add_rect(sl, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.05), fill=C_BLUE)
add_text(sl, "特徴：設計がシンプルで検証しやすく，パイプライン化の基準（ベースライン）として機能",
         Inches(0.7), Inches(6.6), Inches(11.9), Inches(0.45),
         font_size=Pt(14), color=C_LGRAY, italic=True)

# ===================================================
# スライド 7: 高速化設計の詳細（2段パイプライン）
# ===================================================
sl = add_slide()
add_header_bar(sl, "高速化設計の詳細 ② 2段パイプラインCPU", "04 / 2-Stage Pipeline")
page_num(sl, 7)

# 左: 説明
add_card(sl, Inches(0.5), Inches(1.4), Inches(7.2), Inches(5.8))
add_text(sl, "設計ポイント",
         Inches(0.8), Inches(1.5), Inches(6.8), Inches(0.5),
         font_size=Pt(20), bold=True, color=C_CYAN)
lines = [
    ("●", "IF（命令フェッチ）と EX（実行）の2ステージ"),
    ("●", "パイプラインレジスタで段間を分割し\n　クリティカルパスを短縮"),
    ("●", "周期：14 ns → 71.4 MHz（逐次比 +43%）"),
    ("●", "ハザード対策：データハザード時はストール"),
    ("●", "CPI：1.1514（ストールにより増加）"),
]
bullet_lines(sl, lines, Inches(0.9), Inches(2.1), Inches(6.6), Inches(4.5),
             size=Pt(16), spacing=Inches(0.72))

# 右: ステージ図（テキストで）
add_card(sl, Inches(7.9), Inches(1.4), Inches(5.0), Inches(5.8), fill=C_ACCENT)
add_text(sl, "パイプライン構造",
         Inches(8.1), Inches(1.5), Inches(4.6), Inches(0.5),
         font_size=Pt(18), bold=True, color=C_CYAN)
stages = ["Stage 1：IF", "Stage 2：EX / MEM / WB"]
colors_s = [C_BLUE, C_GREEN]
for i, (s, c) in enumerate(zip(stages, colors_s)):
    add_rect(sl, Inches(8.2), Inches(2.2 + i * 2.0), Inches(4.3), Inches(1.5), fill=c)
    add_text(sl, s,
             Inches(8.2), Inches(2.2 + i * 2.0 + 0.5), Inches(4.3), Inches(0.55),
             font_size=Pt(16), bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
    if i < len(stages) - 1:
        add_text(sl, "↓",
                 Inches(9.8), Inches(2.2 + i * 2.0 + 1.5), Inches(0.8), Inches(0.45),
                 font_size=Pt(22), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

# ===================================================
# スライド 8: 高速化設計の詳細（5段パイプライン）
# ===================================================
sl = add_slide()
add_header_bar(sl, "高速化設計の詳細 ③ 5段パイプラインCPU", "04 / 5-Stage Pipeline")
page_num(sl, 8)

add_card(sl, Inches(0.5), Inches(1.4), Inches(7.2), Inches(5.8))
add_text(sl, "設計ポイント",
         Inches(0.8), Inches(1.5), Inches(6.8), Inches(0.5),
         font_size=Pt(20), bold=True, color=C_CYAN)
lines = [
    ("●", "IF → ID → EX → MEM → WB の5ステージ構成"),
    ("●", "各ステージ間にパイプラインレジスタを配置"),
    ("●", "周期：12 ns → 83.3 MHz（逐次比 +67%）"),
    ("●", "フォワーディングによるデータハザード軽減"),
    ("●", "CPI：1.3934（ストール・フラッシュによる増加）"),
    ("●", "LUT使用率：3.65%（小型メモリ構成のため大幅削減）"),
]
bullet_lines(sl, lines, Inches(0.9), Inches(2.1), Inches(6.6), Inches(4.5),
             size=Pt(16), spacing=Inches(0.62))

# 右: 5段ステージ図
add_card(sl, Inches(7.9), Inches(1.4), Inches(5.0), Inches(5.8), fill=C_ACCENT)
add_text(sl, "5段パイプライン構造",
         Inches(8.1), Inches(1.5), Inches(4.6), Inches(0.5),
         font_size=Pt(16), bold=True, color=C_CYAN)
stage5 = ["IF  命令フェッチ", "ID  命令デコード", "EX  演算実行",
          "MEM メモリアクセス", "WB  ライトバック"]
for i, s in enumerate(stage5):
    add_rect(sl, Inches(8.2), Inches(2.1 + i * 0.9), Inches(4.3), Inches(0.75), fill=C_BLUE)
    add_text(sl, s,
             Inches(8.2), Inches(2.1 + i * 0.9 + 0.1), Inches(4.3), Inches(0.55),
             font_size=Pt(14), bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
    if i < 4:
        add_text(sl, "↓", Inches(9.8), Inches(2.1 + i * 0.9 + 0.75), Inches(0.8), Inches(0.35),
                 font_size=Pt(14), color=C_WHITE, align=PP_ALIGN.CENTER)

# ===================================================
# スライド 9: 評価結果（実行時間・CPI・周波数）
# ===================================================
sl = add_slide()
add_header_bar(sl, "評価結果：実行時間・CPI・動作周波数", "05 / Performance Results")
page_num(sl, 9)

# 表
headers = ["実装", "命令数", "動作周波数", "実行クロック数", "CPI", "実行時間", "高速化率"]
col_x = [Inches(0.25), Inches(2.0), Inches(3.6), Inches(5.2), Inches(7.5), Inches(8.8), Inches(10.4)]
col_w2 = [Inches(1.7), Inches(1.55), Inches(1.55), Inches(2.2), Inches(1.25), Inches(1.55), Inches(2.7)]

add_rect(sl, Inches(0.25), Inches(1.45), Inches(12.8), Inches(0.55), fill=C_BLUE)
for j, h in enumerate(headers):
    add_text(sl, h, col_x[j], Inches(1.47), col_w2[j], Inches(0.5),
             font_size=Pt(13), bold=True, color=C_DARK, align=PP_ALIGN.CENTER)

data_rows = [
    ("クイック 逐次",     "118,777", "50 MHz",  "118,778",  "1.000", "2.376 ms", "基準"),
    ("クイック 2段 PL",   "118,777", "71.4 MHz","136,760",  "1.151", "1.914 ms", "1.24×"),
    ("クイック 5段 PL",   "118,777", "83.3 MHz","165,507",  "1.393", "1.986 ms", "1.20×"),
]
row_colors = [C_MID, C_ACCENT, C_MID]
note_colors = [C_LGRAY, C_GREEN, C_CYAN]

for i, (row, rc, nc) in enumerate(zip(data_rows, row_colors, note_colors)):
    add_rect(sl, Inches(0.25), Inches(2.0 + i * 1.35), Inches(12.8), Inches(1.2), fill=rc)
    for j, val in enumerate(row):
        c = nc if j == 6 else C_WHITE
        add_text(sl, val, col_x[j], Inches(2.0 + i * 1.35 + 0.3),
                 col_w2[j], Inches(0.6),
                 font_size=Pt(14), color=c, align=PP_ALIGN.CENTER)

add_text(sl,
         "実行時間 = 実行命令数 × CPI × クロック周期",
         Inches(0.7), Inches(6.15), Inches(11.9), Inches(0.45),
         font_size=Pt(15), bold=True, color=C_YELLOW, align=PP_ALIGN.CENTER)
add_text(sl,
         "2段PL：周波数向上でCPIの増加を上回り最速．5段PLはさらに高周波だが，ハザード増によるCPIの上昇が大きい",
         Inches(0.7), Inches(6.65), Inches(11.9), Inches(0.55),
         font_size=Pt(13), color=C_LGRAY, italic=True)

# ===================================================
# スライド 10: 評価結果（回路資源）
# ===================================================
sl = add_slide()
add_header_bar(sl, "評価結果：回路資源利用状況", "05 / Resource Utilization")
page_num(sl, 10)

headers = ["実装", "Slice LUTs", "（うちLogic）", "（うちMemory）", "Registers", "LUT使用率"]
col_x3 = [Inches(0.3), Inches(2.3), Inches(4.2), Inches(6.3), Inches(8.6), Inches(10.6)]
col_w3 = [Inches(1.9), Inches(1.85), Inches(2.05), Inches(2.25), Inches(1.95), Inches(2.0)]

add_rect(sl, Inches(0.3), Inches(1.45), Inches(12.7), Inches(0.55), fill=C_BLUE)
for j, h in enumerate(headers):
    add_text(sl, h, col_x3[j], Inches(1.47), col_w3[j], Inches(0.5),
             font_size=Pt(13), bold=True, color=C_DARK, align=PP_ALIGN.CENTER)

res_rows = [
    ("クイック 逐次",   "18,957", "2,529",  "16,428", "64",  "29.90%"),
    ("2段 パイプライン","18,888", "2,460",  "16,428", "246", "29.79%"),
    ("5段 パイプライン", "2,316", "1,248",  "1,068",  "657", "3.65%"),
]
row_colors = [C_MID, C_ACCENT, C_MID]
for i, (row, rc) in enumerate(zip(res_rows, row_colors)):
    add_rect(sl, Inches(0.3), Inches(2.05 + i * 1.45), Inches(12.7), Inches(1.3), fill=rc)
    for j, val in enumerate(row):
        c = C_YELLOW if (i == 2 and j in [1, 4]) else C_WHITE
        add_text(sl, val, col_x3[j], Inches(2.05 + i * 1.45 + 0.35),
                 col_w3[j], Inches(0.6),
                 font_size=Pt(14), color=c, align=PP_ALIGN.CENTER)

add_text(sl, "5段パイプラインはメモリサイズを大幅に削減（Distributed RAMを最小限に）",
         Inches(0.7), Inches(6.5), Inches(11.9), Inches(0.45),
         font_size=Pt(16), bold=True, color=C_GREEN)
add_text(sl, "→ LUT使用率が 29.90% → 3.65% へ激減（約1/8）",
         Inches(0.7), Inches(6.95), Inches(11.9), Inches(0.4),
         font_size=Pt(15), color=C_CYAN)

# ===================================================
# スライド 11: 評価結果（消費電力・消費エネルギ）
# ===================================================
sl = add_slide()
add_header_bar(sl, "評価結果：消費電力・消費エネルギ", "05 / Power & Energy")
page_num(sl, 11)

# 左: 電力表
add_card(sl, Inches(0.4), Inches(1.45), Inches(6.2), Inches(5.7))
add_text(sl, "消費電力",
         Inches(0.7), Inches(1.55), Inches(5.5), Inches(0.5),
         font_size=Pt(18), bold=True, color=C_CYAN)

pw_data = [
    ("実装",           "総電力",  "動的電力", "静的電力"),
    ("クイック逐次",   "0.384 W", "0.278 W", "0.105 W"),
    ("2段パイプライン","0.470 W", "0.364 W", "0.105 W"),
    ("5段パイプライン","0.225 W", "0.120 W", "0.105 W"),
]
for i, row in enumerate(pw_data):
    bg = C_BLUE if i == 0 else (C_ACCENT if i % 2 else C_MID)
    add_rect(sl, Inches(0.5), Inches(2.1 + i * 1.0), Inches(6.0), Inches(0.9), fill=bg)
    xs = [Inches(0.55), Inches(2.2), Inches(3.7), Inches(5.0)]
    ws = [Inches(1.6), Inches(1.45), Inches(1.25), Inches(1.3)]
    for j, v in enumerate(row):
        c = C_DARK if i == 0 else (C_YELLOW if (i == 3 and j > 0) else C_WHITE)
        add_text(sl, v, xs[j], Inches(2.1 + i * 1.0 + 0.18),
                 ws[j], Inches(0.55),
                 font_size=Pt(13), bold=(i == 0), color=c)

# 右: 消費エネルギ
add_card(sl, Inches(6.9), Inches(1.45), Inches(6.0), Inches(5.7))
add_text(sl, "消費エネルギ（電力 × 実行時間）",
         Inches(7.1), Inches(1.55), Inches(5.6), Inches(0.5),
         font_size=Pt(16), bold=True, color=C_CYAN)

# Energy: Power × Time
en_data = [
    ("クイック逐次",    0.384, 2.376),   # W, ms
    ("2段 パイプライン",0.470, 1.914),
    ("5段 パイプライン",0.225, 1.986),
]
energies = [(name, p * t, p, t) for name, p, t in en_data]
max_e = max(e[1] for e in energies)

for i, (name, e, p, t) in enumerate(energies):
    bar_w = Inches(3.5 * e / max_e)
    add_rect(sl, Inches(7.1), Inches(2.4 + i * 1.45), bar_w, Inches(0.6),
             fill=C_BLUE if i < 2 else C_GREEN)
    add_text(sl, f"{name}",
             Inches(7.1), Inches(2.1 + i * 1.45), Inches(5.5), Inches(0.35),
             font_size=Pt(13), color=C_LGRAY)
    add_text(sl, f"{e:.3f} mJ",
             Inches(7.1) + bar_w + Inches(0.1), Inches(2.4 + i * 1.45),
             Inches(1.8), Inches(0.55),
             font_size=Pt(14), bold=True,
             color=C_GREEN if i == 2 else C_WHITE)

add_text(sl, "5段PLは消費電力・エネルギ両面で最も優秀",
         Inches(7.1), Inches(6.55), Inches(5.6), Inches(0.45),
         font_size=Pt(15), bold=True, color=C_GREEN)

# ===================================================
# スライド 12: 考察
# ===================================================
sl = add_slide()
add_header_bar(sl, "考察", "06 / Discussion")
page_num(sl, 12)

items_disc = [
    ("実行時間の改善",
     "2段パイプラインは周波数向上（50→71.4MHz）がCPIの増加（1.00→1.15）を上回り\n"
     "実行時間を約20%短縮できた．5段パイプラインは更に高周波（83.3MHz）だが，\n"
     "ハザードによるCPIの大幅増加（1.39）で改善幅は約20%にとどまった．"),
    ("消費電力・エネルギの改善",
     "5段パイプラインはメモリ構成の最適化によりLUT使用率が約1/8となり，\n"
     "動的電力を0.278W→0.120Wへ大幅削減．消費エネルギも逐次比で約51%削減．"),
    ("トレードオフの整理",
     "パイプライン段数を増やすと①動作周波数の向上②CPI増加のトレードオフが生じる．\n"
     "今回の5段PLはメモリ削減の効果が大きく，電力・エネルギ面では最良の結果を示した．"),
]
for i, (title, body) in enumerate(items_disc):
    add_card(sl, Inches(0.5), Inches(1.45 + i * 1.95), Inches(12.3), Inches(1.8))
    add_rect(sl, Inches(0.5), Inches(1.45 + i * 1.95), Inches(0.12), Inches(1.8),
             fill=C_BLUE)
    add_text(sl, title,
             Inches(0.75), Inches(1.5 + i * 1.95), Inches(11.5), Inches(0.45),
             font_size=Pt(16), bold=True, color=C_CYAN)
    add_text(sl, body,
             Inches(0.75), Inches(1.95 + i * 1.95), Inches(11.5), Inches(1.2),
             font_size=Pt(14), color=C_WHITE)

# ===================================================
# スライド 13: まとめ
# ===================================================
sl = add_slide()
add_header_bar(sl, "まとめ", "06 / Summary")
page_num(sl, 13)

summary = [
    ("✓", "クイックソートは選択ソートに比べ実行命令数が約29倍少なくアルゴリズムに採用"),
    ("✓", "2段パイプライン：実行時間を約20%短縮（2.376ms → 1.914ms）"),
    ("✓", "5段パイプライン：消費エネルギを約51%削減（0.912mJ → 0.447mJ）"),
    ("✓", "5段パイプラインはLUT使用率を29.90% → 3.65%へ大幅削減"),
    ("✓", "パイプライン化は周波数向上とCPI増加のトレードオフを伴う"),
]
for i, (mark, text) in enumerate(summary):
    add_rect(sl, Inches(0.5), Inches(1.5 + i * 1.05), Inches(12.3), Inches(0.9),
             fill=C_MID if i % 2 == 0 else C_ACCENT)
    add_text(sl, mark,
             Inches(0.65), Inches(1.5 + i * 1.05 + 0.15), Inches(0.5), Inches(0.6),
             font_size=Pt(18), bold=True, color=C_GREEN)
    add_text(sl, text,
             Inches(1.15), Inches(1.5 + i * 1.05 + 0.15), Inches(11.3), Inches(0.6),
             font_size=Pt(17), color=C_WHITE)

add_text(sl, "→ 目的・用途に応じてアルゴリズム・アーキテクチャを選択することが重要",
         Inches(0.7), Inches(6.85), Inches(11.9), Inches(0.45),
         font_size=Pt(16), bold=True, color=C_YELLOW, align=PP_ALIGN.CENTER)

# ===================================================
# スライド 14: 感想・役割分担・エフォート
# ===================================================
sl = add_slide()
add_header_bar(sl, "感想・役割分担・エフォート", "07 / Impressions & Contributions")
page_num(sl, 14)

# 役割分担テーブル
add_rect(sl, Inches(0.5), Inches(1.45), Inches(12.3), Inches(0.5), fill=C_BLUE)
cols_r = [Inches(0.55), Inches(2.3), Inches(6.2), Inches(10.5)]
cols_rw = [Inches(1.7), Inches(3.85), Inches(4.25), Inches(2.65)]
for j, h in enumerate(["メンバー", "担当内容", "感想", "エフォート"]):
    add_text(sl, h, cols_r[j], Inches(1.47), cols_rw[j], Inches(0.45),
             font_size=Pt(14), bold=True, color=C_DARK)

member_data = [
    ("メンバー A", "逐次CPUの設計・検証\nクイックソート実装", "パイプライン設計の基礎を理解できた", "25%"),
    ("メンバー B", "2段パイプライン設計\nハザード制御", "ストールの実装が最も難しかった", "25%"),
    ("メンバー C", "5段パイプライン設計\nフォワーディング実装", "電力削減の効果が大きく達成感があった", "25%"),
    ("メンバー D", "評価・測定・スライド作成\nレポート執筆", "結果を数値で確認でき理解が深まった", "25%"),
]
for i, (name, role, impression, effort) in enumerate(member_data):
    bg = C_MID if i % 2 == 0 else C_ACCENT
    add_rect(sl, Inches(0.5), Inches(1.95 + i * 1.2), Inches(12.3), Inches(1.1), fill=bg)
    vals = [name, role, impression, effort]
    for j, v in enumerate(vals):
        add_text(sl, v, cols_r[j], Inches(1.95 + i * 1.2 + 0.15),
                 cols_rw[j], Inches(0.8),
                 font_size=Pt(13), color=C_WHITE)

add_text(sl, "※ 各メンバーの氏名・感想・エフォートは実際の内容に書き換えてください",
         Inches(0.7), Inches(7.05), Inches(11.9), Inches(0.35),
         font_size=Pt(12), color=C_LGRAY, italic=True)

# ===================================================
# 保存
# ===================================================
out = "/mnt/c/EX3/Five_pipeline/slides/processor_speedup.pptx"
prs.save(out)
print(f"Saved: {out}")
